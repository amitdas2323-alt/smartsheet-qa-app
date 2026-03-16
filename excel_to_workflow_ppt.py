"""
Read Performance Manager Feasibility Analysis Excel, extract Technology Current State
from yellow-highlighted rows only, and create a workflow-style PowerPoint.
"""

from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE


# Paths
EXCEL_PATH = Path(
    r"c:\Users\amit.das1\OneDrive - JLL\Documents\Agentic AI Requirements"
    r"\Performance Manager Feasibility Analysis_TraciShook.xlsx"
)
OUTPUT_PPTX_PATH = Path(
    r"c:\Users\amit.das1\OneDrive - JLL\Documents\Agentic AI Requirements"
    r"\Performance_Manager_Technology_Current_State_Workflow.pptx"
)
# Fallback if that folder is read-only or locked
OUTPUT_PPTX_FALLBACK = Path(__file__).resolve().parent / "Performance_Manager_Technology_Current_State_Workflow.pptx"

# Yellow detection: Excel yellow highlights (various shades)
YELLOW_RGB_6 = {
    "FFFF00", "FFFF99", "FFFFC7", "FFF2CC", "FFEB9C",
    "FFE699", "FFD966", "FFCC00", "FFFACD", "FFFFE0",
}
# 8-char with alpha
YELLOW_RGB_8 = {s + "00" for s in YELLOW_RGB_6} | {"FFFFFF00", "00FFFF00"}
# Indexed color 64 is often yellow in Excel; 13, 43 can be yellow shades
YELLOW_INDEXED = {64, 13, 43, 44, 45}


def _normalize_hex(s: str | None) -> str:
    if not s or not isinstance(s, str):
        return ""
    s = s.upper().strip()
    if len(s) == 8 and s.startswith("00"):
        s = s[2:]
    if len(s) == 8:
        s = s[2:]  # drop alpha
    return s[:6] if len(s) >= 6 else s


def _is_yellow_color(cell) -> bool:
    """Return True if cell has a yellow fill (fgColor or bgColor)."""
    if not cell or not hasattr(cell, "fill") or not cell.fill:
        return False
    fill = cell.fill
    for color_attr in ("fgColor", "bgColor", "front", "back"):
        color = getattr(fill, color_attr, None)
        if color is None:
            continue
        rgb = getattr(color, "rgb", None)
        if rgb:
            h = _normalize_hex(rgb)
            if h in YELLOW_RGB_6 or h in YELLOW_RGB_8:
                return True
            if "FFFF" in h or h.endswith("FF00") or h == "FFFF00":
                return True
        idx = getattr(color, "indexed", None)
        if idx is not None and idx in YELLOW_INDEXED:
            return True
    return False


def _row_has_yellow_highlight(ws, row_idx: int, max_col: int) -> bool:
    """Return True if any cell in this row has yellow fill."""
    for col in range(1, max_col + 1):
        if _is_yellow_color(ws.cell(row=row_idx, column=col)):
            return True
    return False


def get_technology_current_state_from_excel(excel_path: Path) -> list[str]:
    """Load Excel, find 'Technology Current State' column, return list of values from yellow-highlighted rows only."""
    wb = load_workbook(excel_path, data_only=True)
    ws = wb.active
    max_col = ws.max_column
    max_row = ws.max_row

    # Find header row and column index for "Technology Current State"
    tech_col = None
    header_row = None
    for row in range(1, min(max_row + 1, 30)):
        for col in range(1, max_col + 1):
            val = ws.cell(row=row, column=col).value
            if val and "Technology Current State" in str(val).strip():
                tech_col = col
                header_row = row
                break
        if tech_col is not None:
            break

    if tech_col is None or header_row is None:
        # Try case-insensitive and partial match
        for row in range(1, min(max_row + 1, 30)):
            for col in range(1, max_col + 1):
                val = ws.cell(row=row, column=col).value
                if val and "technology" in str(val).lower() and "current state" in str(val).lower():
                    tech_col = col
                    header_row = row
                    break
            if tech_col is not None:
                break

    if tech_col is None:
        raise ValueError("Column 'Technology Current State' not found in the Excel file.")

    results = []
    for row_idx in range(header_row + 1, max_row + 1):
        if not _row_has_yellow_highlight(ws, row_idx, max_col):
            continue
        val = ws.cell(row=row_idx, column=tech_col).value
        if val is not None and str(val).strip():
            results.append(str(val).strip())
    wb.close()
    return results


def _to_emu(measure) -> int:
    """Convert Length (Inches/Emu) to integer EMU for connectors."""
    if hasattr(measure, "emu"):
        return int(measure.emu)
    return int(measure)


def create_workflow_presentation(items: list[str], output_path: Path) -> None:
    """Create one slide matching Administrative Automation: 3 sections, dashed AI Engine, team→doc→engine→envelope→back to team."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ---- Title (top left, like reference) ----
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = "Administrative Automation"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True

    # Distribute 5 items across 3 sections: 2, 2, 1
    section_items = (
        items[0:2] if len(items) >= 2 else items,
        items[2:4] if len(items) >= 4 else (items[2:3] if len(items) > 2 else []),
        items[4:5] if len(items) >= 5 else [],
    )
    section_titles = (
        ("Performance Manager Assistant", "Agent Flow: Finance Monitoring"),
        ("Work Order / Operations Assistant", "Agent Flow: Corrigo Monitoring"),
        ("ERM / Risk Assistant", "Agent Flow: ERM Monitoring"),
    )
    section_teams = ("Performance Manager Team", "Operations/Engineering Team", "Risk/Operations Team")
    section_data_sources = ("E1 Data", "Corrigo", "EHS, Incidents, Risk, Corrigo")
    section_outputs = ("AI enhanced Financial Report", "AI enhanced Work Order Report", "AI enhanced Risk Report")

    margin_left = Inches(0.45)
    col_team_w = Inches(1.35)
    col_doc_w = Inches(1.2)
    col_engine_w = Inches(4.8)
    col_agents_w = Inches(1.6)
    col_backend_w = Inches(2.0)
    gap = Inches(0.2)
    section_height = Inches(2.0)
    row0 = Inches(1.0)

    # Shared Backend (right side, full height) - one box
    x_backend = margin_left + col_team_w + gap + col_doc_w + gap + col_engine_w + gap + col_agents_w + gap
    backend_top = row0
    backend_h = section_height * 3 + Inches(0.2)
    backend_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x_backend, backend_top, col_backend_w, backend_h,
    )
    backend_box.fill.solid()
    backend_box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    backend_box.line.color.rgb = RGBColor(0x60, 0x60, 0x60)
    backend_txt = slide.shapes.add_textbox(x_backend + Inches(0.15), backend_top + backend_h / 2 - Inches(0.5), col_backend_w - Inches(0.3), Inches(1.0))
    backend_txt.text_frame.word_wrap = True
    backend_txt.text_frame.paragraphs[0].text = "Backend Database / Storage"
    backend_txt.text_frame.paragraphs[0].font.size = Pt(10)
    backend_txt.text_frame.paragraphs[0].font.bold = True
    backend_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    x_engine_left = margin_left + col_team_w + gap + col_doc_w + gap
    x_agents_left = x_engine_left + col_engine_w + Inches(0.08)
    engine_right = x_engine_left + col_engine_w
    black = RGBColor(0x22, 0x22, 0x22)
    light_blue = RGBColor(0xBD, 0xD6, 0xE8)

    # One shared dashed "AI Engine" box spanning all three sections (like reference)
    engine_outer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x_engine_left - Inches(0.02), row0 - Inches(0.02),
        col_engine_w + Inches(0.04), section_height * 3 + Inches(0.04),
    )
    engine_outer.fill.solid()
    engine_outer.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    engine_outer.line.color.rgb = black
    engine_outer.line.width = Pt(0.5)
    engine_outer.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    eng_title = slide.shapes.add_textbox(x_engine_left, row0, col_engine_w, Inches(0.3))
    eng_title.text_frame.paragraphs[0].text = "AI Engine"
    eng_title.text_frame.paragraphs[0].font.bold = True
    eng_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    eng_title.text_frame.paragraphs[0].font.size = Pt(11)

    for sec in range(3):
        sec_items = section_items[sec]
        if not sec_items:
            continue
        top = row0 + sec * section_height
        mid_y = top + section_height / 2
        title_main, title_flow = section_titles[sec]
        team_name = section_teams[sec]
        data_src = section_data_sources[sec]
        out_label = section_outputs[sec]

        # Section titles (left)
        tbox = slide.shapes.add_textbox(margin_left, top - Inches(0.35), Inches(4), Inches(0.28))
        tbox.text_frame.paragraphs[0].text = title_main
        tbox.text_frame.paragraphs[0].font.bold = True
        tbox.text_frame.paragraphs[0].font.size = Pt(12)
        fbox = slide.shapes.add_textbox(margin_left, top - Inches(0.08), Inches(4), Inches(0.25))
        fbox.text_frame.paragraphs[0].text = title_flow
        fbox.text_frame.paragraphs[0].font.size = Pt(10)

        # Left: Team box
        team_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin_left, top + Inches(0.1), col_team_w, Inches(0.5))
        team_box.fill.solid()
        team_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        team_box.line.color.rgb = black
        tmtxt = slide.shapes.add_textbox(margin_left, top + Inches(0.18), col_team_w, Inches(0.35))
        tmtxt.text_frame.paragraphs[0].text = team_name
        tmtxt.text_frame.paragraphs[0].font.size = Pt(8)
        tmtxt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Doc box
        x_doc = margin_left + col_team_w + gap
        doc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_doc, top + Inches(0.1), col_doc_w, Inches(0.5))
        doc_box.fill.solid()
        doc_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        doc_box.line.color.rgb = black
        doctxt = slide.shapes.add_textbox(x_doc, top + Inches(0.18), col_doc_w, Inches(0.35))
        doctxt.text_frame.paragraphs[0].text = data_src
        doctxt.text_frame.paragraphs[0].font.size = Pt(8)
        doctxt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Arrows: Team → Doc → Engine
        arrow_y_emu = _to_emu(mid_y)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(_to_emu(margin_left + col_team_w)), Emu(arrow_y_emu), Emu(_to_emu(x_doc)), Emu(arrow_y_emu)).line.color.rgb = black
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(_to_emu(x_doc + col_doc_w)), Emu(arrow_y_emu), Emu(_to_emu(x_engine_left)), Emu(arrow_y_emu)).line.color.rgb = black

        # Light blue rounded step box inside AI Engine (per section)
        step_box_h = section_height - Inches(0.5)
        step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_engine_left + Inches(0.12), top + Inches(0.32), col_engine_w - Inches(0.24), step_box_h)
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = light_blue
        step_box.line.color.rgb = RGBColor(0x5B, 0x8D, 0xB0)
        # Numbered steps from Excel
        step_font = Pt(8)
        for i, text in enumerate(sec_items):
            y = top + Inches(0.4) + i * (step_box_h / max(len(sec_items), 1)) * 0.85
            st = slide.shapes.add_textbox(x_engine_left + Inches(0.2), y, col_engine_w - Inches(0.5), Inches(0.35))
            st.text_frame.word_wrap = True
            st.text_frame.paragraphs[0].text = f"{i + 1}. {(text[:70] + '...') if len(text) > 70 else text}"
            st.text_frame.paragraphs[0].font.size = step_font
        # Agents to the right of the blue box
        agents = ("Data Retrieval Agent", "Analyst Agent", "Email Agent")
        for ai, agent in enumerate(agents):
            ag = slide.shapes.add_textbox(x_agents_left, top + Inches(0.45) + ai * Inches(0.32), col_agents_w, Inches(0.28))
            ag.text_frame.paragraphs[0].text = agent
            ag.text_frame.paragraphs[0].font.size = Pt(8)

        # Arrow Engine ↔ Backend (double-headed: two straight arrows)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(_to_emu(engine_right)), Emu(arrow_y_emu), Emu(_to_emu(x_backend)), Emu(arrow_y_emu)).line.color.rgb = black
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(_to_emu(x_backend)), Emu(arrow_y_emu), Emu(_to_emu(engine_right)), Emu(arrow_y_emu)).line.color.rgb = black

        # Output: envelope + label (below/left of section)
        env_y = top + section_height - Inches(0.4)
        env_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin_left + Inches(0.2), env_y, Inches(0.55), Inches(0.4))
        env_box.fill.solid()
        env_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        env_box.line.color.rgb = black
        out_txt = slide.shapes.add_textbox(margin_left, env_y + Inches(0.42), col_team_w + col_doc_w, Inches(0.3))
        out_txt.text_frame.paragraphs[0].text = out_label
        out_txt.text_frame.paragraphs[0].font.size = Pt(8)
        # Arrow from engine to envelope
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(_to_emu(x_engine_left + col_engine_w / 2)), Emu(_to_emu(top + section_height - Inches(0.1))),
            Emu(_to_emu(margin_left + col_team_w / 2)), Emu(_to_emu(env_y + Inches(0.2))),
        ).line.color.rgb = black
        # Curved arrow from envelope back to team
        slide.shapes.add_connector(
            MSO_CONNECTOR.CURVE,
            Emu(_to_emu(margin_left + Inches(0.3))), Emu(_to_emu(env_y + Inches(0.2))),
            Emu(_to_emu(margin_left + Inches(0.2))), Emu(_to_emu(top + Inches(0.35))),
        ).line.color.rgb = black

    prs.save(output_path)
    print(f"Saved: {output_path}")


def main() -> None:
    if not EXCEL_PATH.exists():
        print(f"Excel file not found: {EXCEL_PATH}")
        return
    try:
        items = get_technology_current_state_from_excel(EXCEL_PATH)
    except PermissionError:
        print("Could not open the Excel file. Please close it in Excel and run again.")
        return
    if not items:
        print("No yellow-highlighted rows with Technology Current State found.")
        return
    print(f"Found {len(items)} yellow-highlighted row(s) with Technology Current State.")
    output = OUTPUT_PPTX_PATH
    try:
        output.parent.mkdir(parents=True, exist_ok=True)
    except OSError:
        output = OUTPUT_PPTX_FALLBACK
    create_workflow_presentation(items, output)


if __name__ == "__main__":
    main()

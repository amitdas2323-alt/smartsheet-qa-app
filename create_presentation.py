"""
Create a PowerPoint presentation using python-pptx.
Run: pip install -r requirements.txt && python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

def create_presentation(filename="MyPresentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title slide ----
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "My Presentation"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = "Subtitle or date"
    sp.font.size = Pt(24)
    sp.alignment = PP_ALIGN.CENTER

    # ---- Slide 2: Bullet points ----
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Agenda"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = body.text_frame
    points = [
        "Introduction",
        "Key points",
        "Details and data",
        "Conclusion",
        "Questions",
    ]
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(24)
        p.space_after = Pt(12)

    # ---- Slide 3: Content with placeholder ----
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Key Message"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    content.text_frame.word_wrap = True
    p = content.text_frame.paragraphs[0]
    p.text = "Add your main content here. You can edit this file to change titles, bullets, and add more slides."
    p.font.size = Pt(20)

    # ---- Slide 4: Thank you ----
    slide = prs.slides.add_slide(slide_layout)
    thanks = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    p = thanks.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    q = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.6))
    q.text_frame.paragraphs[0].text = "Questions?"
    q.text_frame.paragraphs[0].font.size = Pt(28)
    q.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(filename)
    print(f"Created: {filename}")

if __name__ == "__main__":
    create_presentation()
